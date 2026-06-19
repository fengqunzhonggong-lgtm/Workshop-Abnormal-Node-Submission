"""Generate employee & manager operation guide PPTs."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BLUE = RGBColor(0x1E, 0x40, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)


def make_presentation(title_text, slides_config):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title_text; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "总装车间异常管理系统 V1.0"; p2.font.size = Pt(24); p2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD); p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph(); p3.text = "操作说明"; p3.font.size = Pt(20); p3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE); p3.alignment = PP_ALIGN.CENTER

    # ── Content slides ──
    for idx, (title, bullets) in enumerate(slides_config):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG

        # Top bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
        tf = bar.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"  {idx + 1}. {title}"; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT

        # Step number circle
        circle = slide.shapes.add_shape(9, Inches(0.5), Inches(1.35), Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = WHITE; circle.line.fill.background()
        tf = circle.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.text = str(idx + 1); p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = BLUE; p.alignment = PP_ALIGN.CENTER

        # Content area
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(5))
        tf = txBox.text_frame; tf.word_wrap = True

        for bi, bullet in enumerate(bullets):
            if bi == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = DARK
            p.space_after = Pt(16)
            p.level = 0

    # ── End slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(2))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "感谢使用"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "如有问题请联系系统管理员"; p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD); p2.alignment = PP_ALIGN.CENTER

    return prs


# ═══════════════════════════════════════════════════════════
# Employee PPT
# ═══════════════════════════════════════════════════════════
employee_slides = [
    ("连接系统", [
        "1. 确认手机已连接车间 WiFi 网络",
        "2. 打开微信，扫描车间张贴的二维码",
        "3. 或在微信中直接输入网址：http://[电脑IP]:3000",
        "4. 系统会自动打开登录页面",
    ]),
    ("登录账号", [
        "1. 在登录页面输入您的工号（员工编号）",
        "2. 输入您的登录密码",
        "3. 点击「登录」按钮进入系统",
        "4. 首次登录请联系管理员获取账号密码",
        "5. 登录成功后进入异常填报页面",
    ]),
    ("异常填报", [
        "1. 「工单号」—— 输入对应工单号（必填）",
        "2. 「产品型号」—— 输入或选择产品型号（必填）",
        "3. 「异常类别」—— 从下拉列表选择异常类型（必填）",
        "4. 「数量」—— 输入异常数量，默认为 1",
        "5. 「来源部门」—— 选择异常来源部门（必填）",
        "6. 「来源工序」—— 根据部门自动筛选工序（必填）",
        "7. 「发现工序」—— 选择发现异常的工序（必填）",
        "8. 「异常描述」—— 可填写补充说明（选填）",
        "9. 确认无误后点击「提交」按钮",
        "10. 提交成功后会显示异常编号，请记录下来",
    ]),
    ("查看我的记录", [
        "1. 点击左侧菜单「我的记录」",
        "2. 可查看自己提交的所有异常记录",
        "3. 支持按日期范围和关键词筛选",
        "4. 表格显示编号、工单号、型号、类别、部门、工序、时间等",
        "5. 点击分页按钮可翻页查看更多记录",
    ]),
    ("退出登录", [
        "1. 点击左下角「退出登录」按钮",
        "2. 系统将跳转回登录页面",
        "3. 建议使用完毕后及时退出，保护账号安全",
    ]),
]

# ═══════════════════════════════════════════════════════════
# Manager PPT (includes superadmin)
# ═══════════════════════════════════════════════════════════
manager_slides = [
    ("登录与首页", [
        "1. 确认电脑/手机已连接车间网络",
        "2. 浏览器访问系统地址，或在手机微信扫码",
        "3. 输入管理员账号和密码登录",
        "4. 默认管理员：admin / admin123（请及时修改）",
        "5. 登录后进入异常填报页面",
    ]),
    ("异常填报", [
        "1. 填写工单号、产品型号（可自由输入）、异常类别等信息",
        "2. 选择来源部门后，来源工序会自动联动筛选",
        "3. 选择发现工序，可填写异常描述（选填）",
        "4. 点击「提交」完成填报，系统自动生成异常编号",
        "5. 编号格式：YYYYMMDD + 4位流水号（如202606070001）",
    ]),
    ("异常查询", [
        "1. 点击「异常查询」进入查询页面",
        "2. 可按关键词、工单号、型号、类别、部门、工序、日期等筛选",
        "3. 查询结果以表格形式展示",
        "4. 点击任意一行记录可查看完整详情",
        "5. 支持分页浏览，每页20条记录",
    ]),
    ("Excel 导出", [
        "1. 点击「Excel导出」进入导出页面",
        "2. 可选择开始日期和结束日期限定范围",
        "3. 可选填工单号进行精确导出",
        "4. 点击「导出Excel」按钮下载文件",
        "5. 导出的 Excel 包含所有明细记录",
    ]),
    ("统计分析（管理员）", [
        "1. 点击「统计分析」查看数据图表",
        "2. 顶部可快速选择：今天/本周/本月/今年",
        "3. 也可自定义日期范围进行查询",
        "4. 「部门筛选」下拉可按车间单独分析",
        "5. 展示总异常数、类别分布、部门分布、工序分布、每日趋势",
        "6. 所有图表均为 ECharts 动态渲染",
    ]),
    ("流向分析（管理员）", [
        "1. 点击「流向分析」查看异常流向",
        "2. 支持日期范围筛选",
        "3. 热力图展示来源工序→发现工序的异常流量",
        "4. TOP10 柱状图展示最高频的流向组合",
        "5. 交叉表格显示具体数量分布",
    ]),
    ("基础数据维护（管理员）", [
        "1. 点击「基础数据维护」管理基础信息",
        "2. 可维护：产品型号、异常类别、来源部门、来源工序、发现工序",
        "3. 支持新增、编辑、启用/停用操作",
        "4. 停用的数据不会在填报下拉中显示",
        "5. 产品型号支持在填报时自由输入自动创建",
    ]),
    ("用户管理（管理员）", [
        "1. 点击「用户管理」管理系统账号",
        "2. 可创建新用户：填写工号、姓名、角色、密码",
        "3. 三种角色：系统管理员 / 管理人员 / 员工",
        "4. 支持重置密码、启用/停用账号",
        "5. 停用的账号无法登录系统",
    ]),
    ("系统管理（管理员）", [
        "1. 点击「系统管理」进入管理界面",
        "2. 「运行状态」Tab：查看系统环境检测结果",
        "3. 「备份/归档」Tab：备份数据库、按年份归档数据",
        "4. 「危险操作」Tab：删除指定日期前的历史记录",
        "5. 备份文件保存在程序目录下的 backups 文件夹",
    ]),
    ("环境检测", [
        "1. 程序启动时自动运行10项环境检测",
        "2. 检测项：系统、Python、主机名、网络、磁盘、权限、目录、数据库、文件、端口",
        "3. 系统管理→运行状态可查看完整检测报告",
        "4. 可点击「重新检测」刷新状态",
        "5. 缺失目录会自动创建修复",
    ]),
]


if __name__ == '__main__':
    emp_pptx = make_presentation("员工操作指南", employee_slides)
    emp_pptx.save('release/员工操作指南.pptx')
    print("员工操作指南.pptx saved")

    mgr_pptx = make_presentation("管理人员操作指南", manager_slides)
    mgr_pptx.save('release/管理人员操作指南.pptx')
    print("管理人员操作指南.pptx saved")
